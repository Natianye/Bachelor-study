import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
import os
import re
from pptx.enum.text import PP_ALIGN
import random

def fix_mojibake(text):
    if isinstance(text, str):
        try:
            return text.encode('latin-1').decode('utf-8')
        except Exception:
            try: 
                return text.encode('windows-1252').decode('utf-8')
            except Exception:
                return text
    return text


def simple_markdown_to_text(md):
    if not isinstance(md, str):
        return []

    lines = md.split('\n')
    blocks = []
    buffer = ""
    in_table = False
    in_code_block = False
    table_lines = []
    last_block_type = None

    def flush_buffer(force_newline=False):
        nonlocal buffer, last_block_type
        if buffer.strip():
            blocks.append(buffer.strip())
            if force_newline and blocks and blocks[-1] != "" and not in_code_block:
                blocks.append("")  # add blank line only if not in code block
        buffer = ""
        last_block_type = "paragraph"

    for line in lines:
        stripped = re.sub(r'\s+', ' ', line.strip())

        # Toggle code block state
        if stripped.startswith("```"):
            flush_buffer(force_newline=False)
            blocks.append(stripped)
            in_code_block = not in_code_block
            last_block_type = "code"
            continue

        # Inside code block: keep only non-blank lines
        if in_code_block:
            if line.strip():  # Only add non-blank lines
                blocks.append(line)
            continue

        # Separate trailing ":class:" into its own block
        if ':class:' in stripped and not stripped.startswith(':class:'):
            parts = stripped.rsplit(':class:', 1)
            stripped = parts[0].strip()
            if stripped:
                if buffer:
                    buffer += ' ' + stripped
                else:
                    buffer = stripped
                flush_buffer(force_newline=True)
            blocks.append(":class:")
            last_block_type = "other"
            continue

        # Detect markdown table
        if re.match(r'^\s*\|.*\|\s*', stripped):
            if not in_table:
                flush_buffer(force_newline=True)
                in_table = True
                table_lines = []
            table_lines.append(stripped)
            continue
        else:
            if in_table:
                blocks.append('\n'.join(table_lines))
                in_table = False
                table_lines = []
                blocks.append("")  # Blank line after table

        is_bullet = bool(re.match(r'^\s*([-*+]|[0-9]+\.)\s+', stripped))
        is_heading = bool(re.match(r'^\s*#+\s+', stripped))
        is_logical_break = (
            is_bullet or is_heading or stripped == "" or
            ":class:" in stripped or "(**Note:**" in stripped
        )

        if is_logical_break:
            flush_buffer(force_newline=not is_bullet and not is_heading)
            if stripped:
                blocks.append(stripped)
                last_block_type = (
                    "bullet" if is_bullet else
                    "heading" if is_heading else
                    "other"
                )
        else:
            if buffer:
                buffer += ' ' + stripped
            else:
                buffer = stripped

    if in_table and table_lines:
        blocks.append('\n'.join(table_lines))

    flush_buffer(force_newline=True)

    # Remove consecutive blank lines
    final_blocks = []
    for i, blk in enumerate(blocks):
        if blk == "" and (i == 0 or blocks[i - 1] == ""):
            continue
        final_blocks.append(blk)

    return final_blocks

def unwrap_bullets(md):
    """Combine wrapped bullet lines into single lines, preserve other lines."""
    if not isinstance(md, str):
        return md
    result = []
    buffer = ""
    bullet_re = re.compile(r'^\s*([-*+]|[0-9]+\.)\s+')
    lines = md.split('\n')
    for line in lines:
        if bullet_re.match(line):
            if buffer:
                result.append(buffer)
            buffer = line.strip()
        elif buffer and (line.strip() != ""):
            # Only add to buffer if we're in a bullet and the line is not blank
            buffer += ' ' + line.strip()
        else:
            if buffer:
                result.append(buffer)
                buffer = ""
            result.append(line)
    if buffer:
        result.append(buffer)
    return '\n'.join(result)



def split_text_for_box(paragraphs, first_max=24, other_max=33):
    chunks = []
    current = []
    max_lines = first_max

    for para in paragraphs:
        if len(current) >= max_lines:
            chunks.append(current)
            current = []
            max_lines = other_max  # After the first chunk, switch to other_max
        current.append(para)

    if current:
        chunks.append(current)

    return chunks

# === Load and filter data ===
input_csv = os.path.abspath('./data/both_qa_for_power_point.csv')
df = pd.read_csv(input_csv, encoding='latin-1')
filtered = df[
    (df['Powerpoint'].isin(['7', '8','9'])) &
    (df['Human_Questions'].notnull()) &
    (df['Human_Questions'] != 'No')
].drop_duplicates(subset='chunk_id')

filtered = filtered.head(30)
filtered['content'] = filtered['content'].apply(fix_mojibake)
filtered['Human_Questions'] = filtered['Human_Questions'].apply(fix_mojibake)

prs = Presentation()

# Add cover slide (blank layout)
cover_slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title box with blue background
title_box = cover_slide.shapes.add_shape(
    1,  # Rectangle shape type
    Inches(0.5), Inches(0.7), Inches(9), Inches(1.0)
)
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(220, 235, 255)
title_box.line.fill.background()  # No border

title_frame = title_box.text_frame
title_frame.clear()
p = title_frame.paragraphs[0]
p.text = "Finding the more human-like question?"
p.font.size = Pt(24)         # Slightly smaller font
p.font.bold = True
p.font.underline = False
p.font.name = "Arial"
p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
p.alignment = PP_ALIGN.CENTER

# Main body text box
body_box = cover_slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(5))
body_frame = body_box.text_frame
body_frame.word_wrap = True

# Introduction
p = body_frame.add_paragraph()
p.text = "Introduction:"
p.font.size = Pt(16)
p.font.bold = True
p.font.name = "Arial"

p = body_frame.add_paragraph()
p.text = (
    "Welcome, and thank you for your interest in this survey!\n\n"
    "The goal of my study is to investigate synthetic evaluation QA pairs for RAG systems, using BayBE documentation as the use case. "
    "Specifically, it compares evaluation question-answer (QA) pairs created by LLMs with those created by humans, using the same context."
)

p.font.size = Pt(12)
p.font.name = "Arial"

# Blank line
body_frame.add_paragraph().text = ""

# Procedure
p = body_frame.add_paragraph()
p.text = "Procedure:"
p.font.size = Pt(16)
p.font.bold = True
p.font.name = "Arial"

p = body_frame.add_paragraph()
p.text = (
    "In this survey, you will be shown each page two questions and one context (chunk of BayBE documentation). "
    "Each of the two questions might come from a human or LLM. The two questions are asking for information from the context below them. "
    "Please mark out which one do you think is more 'human-like'?"
)
p.font.size = Pt(12)
p.font.name = "Arial"

# Blank line
body_frame.add_paragraph().text = ""

# Duration
p = body_frame.add_paragraph()
p.text = "Duration:"
p.font.size = Pt(16)
p.font.bold = True
p.font.name = "Arial"

p = body_frame.add_paragraph()
p.text = "This survey will take approximately 15 minutes to complete."
p.font.size = Pt(12)
p.font.name = "Arial"

slide_records = []
pptx_filename = '2nd_experiment_789.pptx'


# === Build slides ===
for idx, row in filtered.iterrows():
    checkbox = u"\u2610"
    human_q = row['Human_Questions']
    gpt_q = row['gpt_41_gs_question']
    unwrapped_content = unwrap_bullets(row['content'])
    context_paragraphs = simple_markdown_to_text(unwrapped_content)
    context_chunks = split_text_for_box(context_paragraphs, first_max=24, other_max=33)

    # Randomly assign which is A and which is B
    if random.choice([True, False]):
        qA, qB = human_q, gpt_q
        labelA, labelB = "Human", "GPT"
    else:
        qA, qB = gpt_q, human_q
        labelA, labelB = "GPT", "Human"
        
    slide_records.append({
        'PowerPoint_Name': pptx_filename,
        'QA_id': row['QA_id'],
        'question_A': labelA,  # "Human" or "GPT"
        'question_B': labelB   # "Human" or "GPT"
    })
   
    
    # First slide with questions and first chunk of context
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Instructions
    box_instr = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.6))
    tf_instr = box_instr.text_frame
    tf_instr.clear()
    tf_instr.word_wrap = True
    tf_instr.vertical_anchor = MSO_ANCHOR.TOP
    p_instr = tf_instr.paragraphs[0]
    p_instr.text = (
        "Here are two questions asking for information from the context below — each of the two questions might come from a human or LLM."
        "Which one do you think is more 'human-like'?\n"
        "(Please replace the checkbox of the 'human-like' question with a 'x')"
    )
    p_instr.font.size = Pt(12)
      
    # Very light blue and very light grey (customize as needed)
    VERY_LIGHT_BLUE = RGBColor(220, 235, 255)
    VERY_LIGHT_GREY = RGBColor(240, 240, 240)    

    # Questions
    boxA = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.6))
    boxA.fill.solid()
    boxA.fill.fore_color.rgb = VERY_LIGHT_BLUE
    tfA = boxA.text_frame
    tfA.clear()
    tfA.word_wrap = True
    tfA.vertical_anchor = MSO_ANCHOR.TOP
    pA = tfA.paragraphs[0] if tfA.paragraphs else tfA.add_paragraph()
    pA.text = f"{checkbox} Question A: {qA}"
    pA.font.size = Pt(12)

    boxB = slide.shapes.add_textbox(Inches(0.4), Inches(1.8), Inches(9.2), Inches(0.6))
    boxB.fill.solid()
    boxB.fill.fore_color.rgb = VERY_LIGHT_BLUE
    tfB = boxB.text_frame
    tfB.clear()
    tfB.word_wrap = True
    tfB.vertical_anchor = MSO_ANCHOR.TOP
    pB = tfB.paragraphs[0] if tfB.paragraphs else tfB.add_paragraph()
    pB.text = f"{checkbox} Question B: {qB}"
    pB.font.size = Pt(12)

    # Context (first chunk, let PowerPoint do the wrapping)
    context_box = slide.shapes.add_textbox(Inches(0.4), Inches(2.5), Inches(9.2), Inches(4.5))
    context_box.fill.solid()
    context_box.fill.fore_color.rgb = VERY_LIGHT_GREY
    tf2 = context_box.text_frame
    tf2.clear()
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.TOP
    
    tf2.text = context_chunks[0][0]  # this creates the first paragraph
    tf2.paragraphs[0].font.size = Pt(11)
    
    for para in context_chunks[0][1:]:
        p = tf2.add_paragraph()
        p.text = para
        p.font.size = Pt(11)
        
    # If there are more context chunks, add the note at the bottom
    if len(context_chunks) > 1:
        p = tf2.add_paragraph()
        p.text = "…more text on next page"
        p.font.size = Pt(10)
        p.font.italic = True

    # Additional slides for more context
    for chunk in context_chunks[1:]:
        if not chunk or all(not para.strip() for para in chunk):
             continue 
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        context_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
        context_box.fill.solid()
        context_box.fill.fore_color.rgb = VERY_LIGHT_GREY
        tf2 = context_box.text_frame
        tf2.clear()
        tf2.word_wrap = True
        if chunk:
            tf2.text = chunk[0]  # Sets the first paragraph
            tf2.paragraphs[0].font.size = Pt(11)
            for para in chunk[1:]:
                p = tf2.add_paragraph()
                p.text = para
                p.font.size = Pt(11)
    

output_df = pd.DataFrame(slide_records, columns=['PowerPoint_Name', 'QA_id', 'question_A', 'question_B'])
output_df.to_csv('2nd_experiment_data.csv', mode='a', index=False, header=False)


prs.save('2nd_experiment_789.pptx')
