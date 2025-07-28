import pandas as pd
from pptx import Presentation
import os

# ==== CONFIGURATION ====
ppt_path = os.path.abspath('./2nd experiment/2nd_experiment_456.pptx')
mapping_csv = '2nd_experiment_data.csv'
questions_csv = './data/both_qa_for_power_point.csv'

# Guess separator for mapping CSV
with open(mapping_csv, 'r', encoding='utf-8') as f:
    first_line = f.readline()
sep = '\t' if '\t' in first_line else ','

# ==== LOAD DATA ====
prs = Presentation(ppt_path)
map_df = pd.read_csv(mapping_csv, sep=sep)
qa_df = pd.read_csv(questions_csv, encoding='latin-1')

# ==== NORMALIZE KEYS ====
map_df['PowerPoint_Name'] = map_df['PowerPoint_Name'].astype(str).str.strip()
qa_df['QA_id'] = qa_df['QA_id'].astype(str).str.strip()
map_df['QA_id'] = map_df['QA_id'].astype(str).str.strip()

ppt_name = os.path.basename(ppt_path).strip().lower()
map_df['PowerPoint_Name'] = map_df['PowerPoint_Name'].str.lower()
pptx_map = map_df[map_df['PowerPoint_Name'] == ppt_name].reset_index(drop=True)

qa_dict = dict(zip(qa_df['QA_id'], qa_df['Human_Questions']))

# ==== UPDATE TITLE ON COVER SLIDE ====
new_title = "Finding the more 'human-like' question?"
cover_slide = prs.slides[0]
title_updated = False
for shape in cover_slide.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.text.strip() != "":
                #print("COVER TITLE BEFORE:", para.text)
                para.text = new_title
                #print("COVER TITLE AFTER:", para.text)
                title_updated = True
                break
    if title_updated:
        break
    
new_procedure_title = "Procedure:"
new_procedure_content = (
    "In this survey, you will be shown on each page two questions and one context (chunk of BayBE documentation). \n"
    "Each question may have been written by either a human or an LLM, and both questions are asking for information from the context below them. \n"
    "Please mark out which one you think is more 'human-like'."
)

# ==== UPDATE PROCEDURE SECTION ON COVER OR INTRO SLIDE ====
# Assume the procedure is on the first or second slide (adjust as needed)
for slide in [prs.slides[0], prs.slides[1]]:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text.strip() == "Procedure:":
                    # print("PROCEDURE TITLE FOUND:", para.text)
                    para.text = new_procedure_title
                    # print("PROCEDURE TITLE UPDATED:", para.text)
                if para.text.strip().startswith("In this survey"):
                    # print("PROCEDURE CONTENT BEFORE:", para.text)
                    para.text = new_procedure_content
                    # print("PROCEDURE CONTENT AFTER:", para.text)

# ==== UPDATE INSTRUCTION TEXT ON QUESTION SLIDES ====
new_instruction = (
        "Here are two questions asking for information from the context below — each of the two questions might come from a human or LLM."
        " Which one do you think is more 'human-like'?\n"
        "(Please replace the checkbox of the 'human-like' question with an 'x')"
)

def is_question_slide(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and ('Question A:' in shape.text or 'Question B:' in shape.text):
            return True
    return False

slide_idx = 0
for idx, row in pptx_map.iterrows():
    # Find the next question slide
    while slide_idx < len(prs.slides) and not is_question_slide(prs.slides[slide_idx]):
        slide_idx += 1
    if slide_idx >= len(prs.slides):
        print("No more slides found for mapping.")
        break

    slide = prs.slides[slide_idx]
    QA_id = row['QA_id']
    human_label = 'A' if row['question_A'] == 'Human' else 'B'
    new_human_q = qa_dict.get(QA_id, None)

    # Update instruction/intro text at top of slide
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if "Here are two questions asking for information from the context below" in para.text:
                para.text = new_instruction

    # Update human question
    if not new_human_q or new_human_q.strip().lower() == 'no':
        print(f"Warning: No updated human question for QA_id {QA_id} in {ppt_name}")
        slide_idx += 1
        continue

    updated = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if f'Question {human_label}:' in para.text:
                prefix = para.text.split('Question')[0]
                para.text = f"{prefix}Question {human_label}: {new_human_q}"
                updated = True
                break
        if updated:
            break

    if not updated:
        print(f"Did not find 'Question {human_label}:' on slide {slide_idx}")

    slide_idx += 1

# ==== SAVE UPDATED PPTX ====
updated_name = ppt_path.replace('.pptx', '_updated.pptx')
prs.save(updated_name)
print(f"\nUpdated PowerPoint saved as {updated_name}")
